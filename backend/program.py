import win32com.client
import pandas as pd
import os
from pptx import Presentation, parts
from pptx.enum.shapes import MSO_SHAPE_TYPE
from io import BytesIO
from pptx.parts.image import Image
import requests
import gdown
import re

def open_presentation(link_presentation):
    # Author: @oceantran27
    # Edit: @thnhmai06
    try:
        ppt_instance = win32com.client.Dispatch("PowerPoint.Application")
        read_only = True
        has_title = False
        window = False
        if ppt_instance:
            prs = ppt_instance.Presentations.open(
                link_presentation, read_only, has_title, window
            )
            return prs, ppt_instance
        return None, None
    except Exception as e:
        print(f"*ERROR - opening presentation: {e}")
        return None, None


def close_presentation(prs, ppt_instance):
    # Author: @oceantran27
    # Edit: @thnhmai06
    try:
        if prs and ppt_instance:
            prs.Close()
            ppt_instance.Quit()
            del ppt_instance
            return True
        return False
    except Exception as e:
        print(f"*ERROR - closing presentation: {e}")
        return False


def duplicate_slide(prs, number_of_copies, slide_index=1):  # count from 1 for win32COM
    # Author: @oceantran27
    # Edit: @thnhmai06

    try:
        if prs:
            for i in range(number_of_copies):
                prs.Slides(slide_index).Duplicate()
            prs.Save()
            return True
        return False
    except Exception as e:
        print(f"*ERROR - duplicating slide: {e}")
        return False


def replace_text_placeholders(prs, data):
    # Author: @oceantran27
    # Edit: @thnhmai06

    column_names = data.columns.tolist()
    for slide_index in range(1, prs.Slides.Count):
        row = data.iloc[slide_index - 1]
        for shape in prs.Slides(slide_index).Shapes:
            if shape.HasTextFrame and shape.TextFrame.HasText:
                text = shape.TextFrame.TextRange.Text
                for column_name in column_names:
                    if column_name in text:
                        new_text = text.replace(column_name, str(row[column_name]))
                        shape.TextFrame.TextRange.Text = new_text
    prs.Save()

def get_image_shape_indices(slide):
    image_indices = []
    for index in range(1, slide.Shapes.Count + 1):
        shape = slide.Shapes(index)
        if shape.Type == 13:
            image_indices.append(index - 1) # count from 0 for pptx
    return image_indices

def save_images_from_shapes(prs_path, images_output_path, shape_indices, slide_index = 0): # count from 0 for pptx
    presentation = Presentation(prs_path)

    if slide_index >= len(presentation.slides):
        raise IndexError("Slide index out of range")

    slide = presentation.slides[slide_index]

    if not os.path.exists(images_output_path):
        os.makedirs(images_output_path)

    for shape_index in shape_indices:
        if shape_index >= len(slide.shapes):
            print(f"Shape index {shape_index} out of range in slide {slide_index}")
            continue

        shape = slide.shapes[shape_index]
        
        if shape.shape_type == 13:
            image = shape.image
            image_bytes = image.blob

            image_path = os.path.join(images_output_path, f"{shape_index + 1}.{image.ext}") # count from 1 for win32COM
            with open(image_path, "wb") as img_file:
                img_file.write(image_bytes)
            print(f"Saved image: {image_path}")
        else:
            print(f"Shape at index {shape_index} is not an image.")

def get_image_extension(url):
    response = requests.head(url, allow_redirects=True)
    content_type = response.headers.get('Content-Type', '')
    extension = ''
    if content_type == 'image/jpeg':
        extension = '.jpg'
    elif content_type == 'image/png':
        extension = '.png'
    elif content_type == 'image/webp':
        extension = '.webp'
    elif content_type == 'image/jpeg':
        extension = '.jpeg'
    return extension

def get_drive_file_id(url):
    match = re.search(r"/file/d/([a-zA-Z0-9_-]+)", url)
    if match:
        return match.group(1)
    else:
        return None

def download_image_from_drive(url, save_path):
    formatted_url = "https://drive.google.com/uc?id="
    try:
        url = formatted_url + get_drive_file_id(url)
        extension = get_image_extension(url)
        if (extension):
            file_path = os.path.join(save_path, f"replace_image{extension}")
            gdown.download(url, file_path, quiet=False)
            print(f"Image has been downloaded and saved to: {file_path}")
            return file_path
        return None

    except requests.exceptions.RequestException as e:
        print(f"Error downloading image: {e}")
        return None

def main():
    # const
    BASE_PATH = "PATH_TO_PROJECT/"
    PRS_PATH = BASE_PATH + "template/template.pptx"
    IMAGES_OUTPUT_PATH = BASE_PATH + "images/template_images"
    DATA_PATH = BASE_PATH + "data/data.xlsx"
    DATA_IMAGE_PATH = BASE_PATH + "images/data_images"

    # init
    data = pd.read_excel(DATA_PATH)
    data_length = data.shape[0]
    prs, ppt_instance = open_presentation(PRS_PATH)

    # process
    if (prs):
        is_duplicate_success = duplicate_slide(prs, data_length, 1)
        if is_duplicate_success:
            replace_text_placeholders(prs, data)
        shape_indices = get_image_shape_indices(prs.Slides(1))
        print(type(prs.Slides(1)))
        save_images_from_shapes(PRS_PATH, IMAGES_OUTPUT_PATH , shape_indices)
        download_data_images(DATA_IMAGE_PATH, data)

    # close
    if (prs and ppt_instance):
        close_presentation(prs, ppt_instance)


if __name__ == "__main__":
    main()
