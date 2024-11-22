from typing import TYPE_CHECKING
if TYPE_CHECKING:
    # Anti-circular import
    from ui.menu import Ui
from globals import input, SHAPES_PATH
import os
from pptx.presentation import Presentation
from pptx.shapes.picture import Picture
from logger.info import console_info
import pandas as pd

def __get_save_path(ui: 'Ui') -> str:
    save_path = ui.save_path.text()
    input.save_path.set(save_path)

def __get_config(ui: 'Ui') -> None:
    def _get_text_config() -> None:
        text_list = ui.config_text_list
        for index in range(text_list):
            input.config.add_text(text_list.item(index).text())

    def _get_image_config() -> None:
        image_table = ui.config_image_table
        for row in range(image_table.rowCount()):
            shape_id = image_table.item(row, 0).text()
            placeholder = image_table.item(row, 1).text()
            input.config.add_image(shape_id=shape_id, placeholder=placeholder)

    _get_text_config()
    _get_image_config()

def get_shapes(prs: Presentation, slide_index = 0, save_path: str = SHAPES_PATH): # Slide đầu tiên có index = 0
    # Author: @oceantran27
    # Edit: @thnhmai06
    # Description: Hàm này sẽ lưu lại các Shapes ảnh (đã xác định trong shape_indices) vào thư mục SHAPES_PATH
    # Edit note: Đã gộp hàm get_image_shape_indices và save_images_from_shapes thành hàm này
    
    IMAGE_TYPE = 13 #ID của shape ảnh trong PowerPoint

    # Tạo folder nếu thư mục lưu không tồn tại
    if not os.path.exists(save_path):
        os.makedirs(save_path)
    # Xóa hết các file trong save_path
    for filename in os.listdir(save_path):
        file_path = os.path.join(save_path, filename)
        if os.path.isfile(file_path):
            os.remove(file_path)

    slide = prs.slides[slide_index]
    for __shape_index_win32COM in range(1, len(slide.shapes) + 1): 
        #__shape_index_win32COM là chỉ số của shape trong slide (theo Win32COM, vì win32COM đếm từ 1)
        # Phần range cộng thêm 1 vì range(a,b) chỉ lấy từ a -> b-1
        
        __shape_index_python_pptx = __shape_index_win32COM - 1 
        # Chỉ số của shape trong slide (theo python-pptx, vì python-pptx đếm từ 0)

        shape = slide.shapes[__shape_index_python_pptx]
        if shape.shape_type == IMAGE_TYPE:
            # Xác nhận rằng shape có kiểu Picture. Comment: Code cháy wá 🔥🔥🔥
            assert isinstance(shape, Picture)

            # Lấy dữ liệu ảnh từ shape
            image = shape.image
            image_bytes = image.blob

            # Lưu ảnh vào thư mục save_path
            image_path = os.path.join(save_path, f"{__shape_index_python_pptx + 1}.{image.ext}")
            with open(image_path, "wb") as img_file:
                img_file.write(image_bytes)
                # Lưu thông tin ảnh vào input.shapes
                input.shapes.add(__shape_index_python_pptx, image_path)
            console_info(__name__, f"Image ID: {__shape_index_win32COM} -> {image_path} (Preview)")

def get_csv(csv_path: str) -> bool:
    '''
    Return: 
    - True: Saved successfully
    - False: CSV is not valid
    '''
    __df = pd.read_csv(csv_path)
    number_of_students = len(__df)
    if not number_of_students>=1:
        return False
    
    input.csv.placeholders = __df.columns.tolist()
    console_info(__name__, "Fields:", (" - ").join(input.csv.placeholders))
    input.csv.students = __df.to_dict(orient='records')
    console_info(__name__, "Students:", f"({len(input.csv.students)})")
    return True

def get_save_path_and_config(ui: 'Ui') -> None:
    __get_save_path(ui)
    __get_config(ui)