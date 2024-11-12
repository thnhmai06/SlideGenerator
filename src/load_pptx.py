import os
from typing import TYPE_CHECKING
from globals import SHAPES_PATH   
from pptx.shapes.picture import Picture
from pptx.presentation import Presentation
from pptx import Presentation as init_presentation
from logger.info import console_info, default as info

if TYPE_CHECKING:
    # Anti-circular import
    from ui.menu import Ui

def __toggle_config_image(ui: 'Ui', is_enable: bool):
    # Enable the config_image_table, add_button, and remove_button
    config_image_table = ui.config_image_table
    add_button = ui.config_image_add_button
    remove_button = ui.config_image_remove_button

    config_image_table.setEnabled(is_enable)
    add_button.setEnabled(is_enable)
    remove_button.setEnabled(is_enable)

def __save_image_shapes_preview(prs: Presentation, slide_index = 0, save_path: str = SHAPES_PATH): # Slide đầu tiên có index = 0
    # Author: @oceantran27
    # Edit: @thnhmai06
    # Description: Hàm này sẽ lưu lại các Shapes ảnh (đã xác định trong shape_indices) vào thư mục SHAPES_PATH
    # Edit note: Đã gộp hàm get_image_shape_indices và save_images_from_shapes thành hàm này
    
    IMAGE_TYPE = 13 #ID của shape ảnh trong PowerPoint

    # Tạo folder nếu thư mục lưu không tồn tại
    if not os.path.exists(save_path):
        os.makedirs(save_path)
    # Xóa hết các file trong save_path
    for filename in os.listdir(save_path):
        file_path = os.path.join(save_path, filename)
        if os.path.isfile(file_path):
            os.remove(file_path)

    slide = prs.slides[slide_index]
    for __shape_index_win32COM in range(1, len(slide.shapes) + 1): 
        #__shape_index_win32COM là chỉ số của shape trong slide (theo Win32COM, vì win32COM đếm từ 1)
        # Phần range cộng thêm 1 vì range(a,b) chỉ lấy từ a -> b-1
        
        __shape_index_python_pptx = __shape_index_win32COM - 1 
        # Chỉ số của shape trong slide (theo python-pptx, vì python-pptx đếm từ 0)

        shape = slide.shapes[__shape_index_python_pptx]
        if shape.shape_type == IMAGE_TYPE:
            # Xác nhận rằng shape có kiểu Picture. Comment: Code cháy wá 🔥🔥🔥
            assert isinstance(shape, Picture)

            # Lấy dữ liệu ảnh từ shape
            image = shape.image
            image_bytes = image.blob

            # Lưu ảnh vào thư mục save_path
            image_path = os.path.join(save_path, f"{__shape_index_python_pptx + 1}.{image.ext}")
            with open(image_path, "wb") as img_file:
                img_file.write(image_bytes)
            console_info(__name__, f"Image ID: {__shape_index_win32COM} -> {image_path} (Preview)")

def __load_preview_to_items_of_config_image_table(ui: 'Ui'):
    config_image_table = ui.config_image_table
    

def load(ui: 'Ui'):
    pptx_path = ui.pptx_path.text()
    prs = init_presentation(pptx_path)

    __toggle_config_image(ui, False)
    ui.config_image_table.clear()

    # Kiểm tra xem trong prs có slide nào không
    if not prs.slides:
        ui.pptx_path.clear() # Xóa đường dẫn file pptx
        info(__name__, "no_slide_pptx")
        return
    
    __save_image_shapes_preview(prs) # Lưu các ảnh từ slide đầu tiên vào thư mục SHAPES_PATH
