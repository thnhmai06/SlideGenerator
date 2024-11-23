import os
from typing import TYPE_CHECKING
from globals import input
from pptx import Presentation
from logger.info import console_info, default as info
from logger.debug import console_debug
from src.toggle_config import toggle_config_image
from src.get_input import get_shapes
from globals import SHAPES_PATH

if TYPE_CHECKING:
    # Anti-circular import
    from ui.menu import Ui

def __refresh_placeholders():
    # Làm mới placeholders ở local file này
    global __placeholders
    __placeholders = input.csv.placeholders

def load(ui: 'Ui'):
    pptx_path = ui.pptx_path.text()
    prs = Presentation(pptx_path)

    toggle_config_image(ui, False)
    ui.config_image_table.clearContents()

    # Nếu prs không có slide nào
    if not prs.slides:
        ui.pptx_path.clear() # Xóa đường dẫn file pptx
        info(__name__, "no_slide_pptx")
        return
    
    # Xóa các ảnh cũ trong thư mục SHAPES_PATH
    if os.path.exists(SHAPES_PATH):
        for filename in os.listdir(SHAPES_PATH):
            file_path = os.path.join(SHAPES_PATH, filename)
            if os.path.isfile(file_path) or os.path.islink(file_path):
                os.unlink(file_path)
            elif os.path.isdir(file_path):
                os.rmdir(file_path)
    
    get_shapes(prs) # Lưu các ảnh từ slide đầu tiên vào thư mục SHAPES_PATH

    ui.config_image_preview.setEnabled(True) # Enable preview button
    # Chỉ khi đã có sẵn placeholder rồi thì mới enable config_image_table
    __refresh_placeholders()
    if (len(__placeholders) > 0):
        toggle_config_image(ui, True)
        