from globals import input, SHAPES_PATH
import os
import io
from PyQt5.QtWidgets import QLineEdit
from pptx.presentation import Presentation
from pptx.shapes.picture import Picture
from src.logger.info import console_info
import polars as pl
from PIL import Image

def _reduce_image_quality(image_bytes: bytes, quality: int) -> bytes:
    """
    Reduce the quality of the image to the minimum visible level.
    
    Args:
    - image_bytes: The original image in bytes.
    - quality: The quality of the reduced image (0-100).
    
    Returns:
    - The reduced quality image in bytes.
    """
    # Open the image from bytes
    image = Image.open(io.BytesIO(image_bytes))
    
    # Create a BytesIO object to save the reduced quality image
    output = io.BytesIO()
    
    # Save the image with the lowest quality
    image.save(output, format = image.format, quality = quality)
    
    # Get the reduced quality image bytes
    reduced_image_bytes = output.getvalue()
    
    return reduced_image_bytes

def get_pptx_path(line_widget: QLineEdit) -> str:
    pptx_path = line_widget.text()
    input.pptx.setPath(pptx_path)
    return pptx_path

def get_csv(csv_path: str) -> bool:
    """
    Return:
    - True: Saved successfully
    - False: CSV is not valid
    """
    LINES_PER_BATCH = 1

    input.csv.df = pl.read_csv(csv_path, batch_size=LINES_PER_BATCH)
    input.csv.placeholders = input.csv.df.columns
    input.csv.number_of_students = len(input.csv.df)

    if not input.csv.number_of_students >= 1:
        return False
    
    console_info(__name__, "Fields:", (" - ").join(input.csv.placeholders))
    console_info(__name__, "Students:", f"({input.csv.number_of_students})")
    return True

def get_shapes(prs: Presentation, slide_index=0, shapes_path: str = SHAPES_PATH):  # Slide đầu tiên có index = 0
    # Author: @oceantran27
    # Edit: @thnhmai06
    # Description: Hàm này sẽ lưu lại các Shapes ảnh (đã xác định trong shape_indices) vào thư mục SHAPES_PATH
    # Edit note: Đã gộp hàm get_image_shape_indices và save_images_from_shapes thành hàm này

    IMAGE_TYPE = 13  # ID của shape ảnh trong PowerPoint
    IMAGE_QUALITY = 5  # Chất lượng ảnh sau khi lưu (0-100)

    # Tạo folder nếu thư mục lưu không tồn tại
    if not os.path.exists(shapes_path):
        os.makedirs(shapes_path)
    # Xóa hết các file trong save_path
    for filename in os.listdir(shapes_path):
        file_path = os.path.join(shapes_path, filename)
        if os.path.isfile(file_path):
            os.remove(file_path)

    slide = prs.slides[slide_index]
    for __shape_index_win32COM in range(1, len(slide.shapes) + 1):
        # __shape_index_win32COM là chỉ số của shape trong slide (theo Win32COM, vì win32COM đếm từ 1)
        # Phần range cộng thêm 1 vì range(a,b) chỉ lấy từ a -> b-1

        __shape_index_python_pptx = __shape_index_win32COM - 1
        # Chỉ số của shape trong slide (theo python-pptx, vì python-pptx đếm từ 0)

        shape = slide.shapes[__shape_index_python_pptx]
        if shape.shape_type == IMAGE_TYPE:
            # Xác nhận rằng shape có kiểu Picture. Comment: Code cháy wá 🔥🔥🔥
            assert isinstance(shape, Picture)

            # Lấy dữ liệu ảnh từ shape
            image = shape.image
            image_bytes = _reduce_image_quality(image.blob, IMAGE_QUALITY)

            # Lưu ảnh vào thư mục save_path
            image_path = os.path.join(
                shapes_path, f"{__shape_index_python_pptx + 1}.{image.ext}"
            )
            with open(image_path, "wb") as img_file:
                img_file.write(image_bytes)
                # Lưu thông tin ảnh vào input.shapes
                input.shapes.add(__shape_index_python_pptx, image_path)
            console_info(
                __name__,
                f"Shape ID: {__shape_index_win32COM} -> {image_path}",
            )

def get_save_path(line_widget: QLineEdit) -> str:
    save_path = line_widget.text()
    input.save.setPath(save_path)
    return save_path