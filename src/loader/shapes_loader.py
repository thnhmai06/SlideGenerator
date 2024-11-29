import os
from typing import TYPE_CHECKING
from globals import input
from pptx import Presentation
from src.logger.info import default as info
from src.loader._toggle_config import toggle_config_image
from src.loader._get_utils import get_shapes, get_pptx_path
from globals import SHAPES_PATH

if TYPE_CHECKING:
    # Anti-circular import
    from src.ui.menu import Menu


def __refresh_placeholders():
    # Làm mới placeholders ở local file này
    global __placeholders
    __placeholders = input.csv.placeholders


def __no_slide(menu: "Menu"):
    menu.pptx_path.clear()  # Xóa đường dẫn file pptx
    info(__name__, "pptx.no_slide")

def __too_much_slide(menu: "Menu"):
    menu.pptx_path.clear()  # Xóa đường dẫn file pptx
    info(__name__, "pptx.too_much_slides")

def __delete_all_file(PATH: str):
    if os.path.exists(PATH):
        for filename in os.listdir(PATH):
            file_path = os.path.join(PATH, filename)
            if os.path.isfile(file_path) or os.path.islink(file_path):
                os.unlink(file_path)
            elif os.path.isdir(file_path):
                os.rmdir(file_path)



def load_shapes(menu: "Menu"):
    get_pptx_path(menu.pptx_path)
    prs = Presentation(input.pptx.path)

    toggle_config_image(menu, False)
    menu.config_image_table.clearContents()

    # Nếu prs không có slide nào
    if not prs.slides:
        __no_slide(menu)    
        return
    if len(prs.slides) > 1:
        __too_much_slide(menu)
        return

    # Xóa các ảnh cũ trong thư mục SHAPES_PATH
    __delete_all_file(SHAPES_PATH)

    # Lưu các ảnh từ slide đầu tiên (0) vào thư mục SHAPES_PATH
    get_shapes(prs, 0) 

    # Enable preview button
    menu.config_image_preview.setEnabled(True)
    
    # Chỉ khi đã có sẵn placeholder rồi thì mới enable config_image_table
    __refresh_placeholders()
    if len(__placeholders) > 0:
        toggle_config_image(menu, True)

    del prs