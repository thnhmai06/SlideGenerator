import os
from typing import TYPE_CHECKING
from globals import user_input
from pptx.presentation import Presentation
from pptx import Presentation as initPresentation
from src.logger.info import default as info
from src.loader._ui_config_controls import toggle_config_image
from globals import SHAPES_PATH
from src.logger.info import console_info
from pptx.shapes.picture import Picture
from PyQt5.QtWidgets import QLineEdit
from src.loader._utils import delete_all_file

if TYPE_CHECKING:
    # Anti-circular import
    from src.ui.menu import Menu


def get_pptx_path(line_widget: QLineEdit) -> str:
    pptx_path = line_widget.text()
    user_input.pptx.setPath(pptx_path)
    return pptx_path

def __load_shapes(
    prs: Presentation, slide_index=0, shapes_path: str = SHAPES_PATH
):  # Slide đầu tiên có index = 0
    # Author: @oceantran27
    # Edit: @thnhmai06
    # Description: Hàm này sẽ lưu lại các Shapes ảnh (đã xác định trong shape_indices) vào thư mục SHAPES_PATH
    # Edit note: Đã gộp hàm get_image_shape_indices và save_images_from_shapes thành hàm này

    IMAGE_TYPE = 13  # ID của shape ảnh trong PowerPoint
    IMAGE_QUALITY = 5  # Chất lượng ảnh sau khi lưu (0-100)  # noqa: F841

    # Tạo folder nếu thư mục lưu không tồn tại
    if not os.path.exists(shapes_path):
        os.makedirs(shapes_path)

    # Xóa hết các file trong save_path
    delete_all_file(shapes_path)

    slide = prs.slides[slide_index]
    for shape_index_win32COM in range(1, len(slide.shapes) + 1):
        # shape_index_win32COM là chỉ số của shape trong slide (theo Win32COM, vì win32COM đếm từ 1) -> Dùng để định danh shape
        # shape_index_python_pptx là chỉ số của shape trong slide (theo python-pptx, vì python-pptx đếm từ 0) -> Dùng để lấy shape (do đang dùng python-pptx)
        #* Ưu tiên đếm theo shape_index_win32COM
        shape_index_python_pptx = shape_index_win32COM - 1

        shape = slide.shapes[shape_index_python_pptx]
        if shape.shape_type == IMAGE_TYPE:
            # Xác nhận rằng shape có kiểu Picture. Comment: Code cháy wá 🔥🔥🔥
            assert isinstance(shape, Picture)

            # Lấy dữ liệu ảnh từ shape
            image = shape.image
            # image_bytes = _reduce_image_quality(image.blob, IMAGE_QUALITY)
            image_bytes = image.blob

            # Lưu ảnh vào thư mục save_path
            image_path = os.path.join(shapes_path, f"{shape_index_win32COM}.{image.ext}")
            with open(image_path, "wb") as img_file:
                img_file.write(image_bytes)
                # Lưu thông tin ảnh vào input.shapes
                user_input.shapes.add(shape_index_win32COM, image_path)
            console_info(__name__, f"Shape ID: {shape_index_win32COM} -> {image_path}")

def __notice_no_slide(menu: "Menu"):
    menu.pptx_path.clear()  # Xóa đường dẫn file pptx
    info(__name__, "pptx.no_slide")

def __notice_too_much_slide(menu: "Menu"):
    menu.pptx_path.clear()  # Xóa đường dẫn file pptx
    info(__name__, "pptx.too_much_slides")


def process_shapes(menu: "Menu"):
    get_pptx_path(menu.pptx_path)
    prs = initPresentation(user_input.pptx.path)

    toggle_config_image(menu, False)
    menu.config_image_table.clearContents()

    # Nếu prs không có slide nào
    if not prs.slides:
        __notice_no_slide(menu)
        return
    if len(prs.slides) > 1:
        __notice_too_much_slide(menu)
        return

    # Xóa các ảnh cũ trong thư mục SHAPES_PATH
    delete_all_file(SHAPES_PATH)

    # Lưu các ảnh từ slide đầu tiên (0) vào thư mục SHAPES_PATH
    __load_shapes(prs, 0)

    # Enable preview button
    menu.config_image_preview.setEnabled(True)

    # Chỉ khi đã có sẵn placeholder rồi thì mới enable config_image_table
    if user_input.csv.number_of_students > 0:
        toggle_config_image(menu, True)
